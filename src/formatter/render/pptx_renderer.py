"""
PPTX Renderer — converts FormattedDocument IR to PowerPoint (.pptx) bytes.

Uses python-pptx to produce a presentation with one slide per page.
Text paragraphs are rendered into text-box shapes with per-run formatting
(bold, italic, underline, font, size, color). Tables are rendered as
native PPTX table shapes with header styling.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from src.formatter.extractor import (
    FormattedDocument,
    FormattedPage,
    FormattedParagraph,
    FormattedSpan,
    FormattedTable,
    FormattedTableCell,
)

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

# Default slide dimensions
_SLIDE_WIDTH_IN = 10.0
_SLIDE_HEIGHT_IN = 7.5
_MARGIN_IN = 0.5
_TABLE_FONT_SIZE = 9


class PPTXRenderer:
    """Renders a FormattedDocument to PPTX bytes.

    Usage::

        renderer = PPTXRenderer()
        pptx_bytes = renderer.render(doc)
    """

    def __init__(
        self,
        default_font: str = "Calibri",
        default_size: float = 11.0,
    ):
        """Initialise the renderer.

        Args:
            default_font: Fallback font family for body text.
            default_size: Fallback font size (pt).
        """
        self._default_font = default_font
        self._default_size = default_size

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, doc: FormattedDocument) -> bytes:
        """Render a FormattedDocument as PPTX bytes.

        Args:
            doc: A FormattedDocument IR instance.

        Returns:
            PPTX file content as bytes.

        Raises:
            ImportError: If python-pptx is not installed.
        """
        if not _HAS_PPTX:
            raise ImportError(
                "python-pptx is required for PPTX rendering. "
                "Install it with: pip install python-pptx"
            )

        prs = Presentation()

        # Set slide dimensions
        prs.slide_width = Inches(_SLIDE_WIDTH_IN)
        prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

        if not doc.pages:
            # Add a single blank slide for an empty document
            layout = prs.slide_layouts[6]  # blank layout
            prs.slides.add_slide(layout)
            return self._save(prs)

        for page in doc.pages:
            self._add_slide(prs, page)

        return self._save(prs)

    # ------------------------------------------------------------------
    # Internal — slide construction
    # ------------------------------------------------------------------

    def _add_slide(self, prs: Any, page: FormattedPage) -> None:
        """Add a slide for a single FormattedPage."""
        # Use blank layout (index 6)
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        content_left = Inches(_MARGIN_IN)
        content_top = Inches(_MARGIN_IN)
        content_width = Inches(_SLIDE_WIDTH_IN - 2 * _MARGIN_IN)
        y_cursor = content_top

        # Render paragraphs and tables sorted by y-position
        items: list[tuple[float, str, Any]] = []  # (y_pos, kind, obj)

        for para in page.paragraphs:
            y_pos = para.lines[0].y_center if para.lines else 0.0
            items.append((y_pos, "para", para))

        for table in page.tables:
            items.append((table.y_position, "table", table))

        items.sort(key=lambda x: x[0])

        for _, kind, obj in items:
            if kind == "para":
                y_cursor = self._add_paragraph_shape(
                    slide, obj, content_left, y_cursor, content_width,
                )
            elif kind == "table":
                y_cursor = self._add_table_shape(
                    slide, obj, content_left, y_cursor, content_width,
                )

    # ------------------------------------------------------------------
    # Internal — paragraphs
    # ------------------------------------------------------------------

    def _add_paragraph_shape(
        self,
        slide: Any,
        para: FormattedParagraph,
        left: Any,
        top: Any,
        width: Any,
    ) -> Any:
        """Add a text box for a paragraph. Returns updated y position."""
        if not para.lines:
            return top

        # Skip image paragraphs (not supported in this simple renderer)
        if para.style == "image":
            return top

        # Estimate height from number of lines
        line_count = len(para.lines)
        font_size = para.font_size
        estimated_height = Pt(font_size * line_count * 1.5 + 10)

        txbox = slide.shapes.add_textbox(left, top, width, estimated_height)
        tf = txbox.text_frame
        tf.word_wrap = True

        # First paragraph in the text frame is created automatically
        pptx_para = tf.paragraphs[0]
        self._set_alignment(pptx_para, para.alignment)

        # Add runs for each span across all lines
        first_run = True
        for line_idx, line in enumerate(para.lines):
            if line_idx > 0:
                # Add a line break by adding a new paragraph
                pptx_para = tf.add_paragraph()
                self._set_alignment(pptx_para, para.alignment)

            for span in line.spans:
                if not span.text:
                    continue

                # List prefix for first span of first line
                text = span.text
                if first_run:
                    if para.style == "list_bullet":
                        text = "\u2022 " + text
                    elif para.style == "list_number":
                        text = "1. " + text

                if first_run:
                    run = pptx_para.runs[0] if pptx_para.runs else pptx_para.add_run()
                    run.text = text
                    first_run = False
                else:
                    run = pptx_para.add_run()
                    run.text = text

                self._apply_run_formatting(run, span, para.style)

        return top + estimated_height

    def _apply_run_formatting(self, run: Any, span: FormattedSpan, style: str) -> None:
        """Apply formatting from a FormattedSpan to a pptx Run."""
        font = run.font

        # Font name
        font_name = span.font or self._default_font
        # Clean up PDF-style font names
        for suffix in ("-Bold", "-Italic", "-BoldItalic", "-Regular", "MT", "PS"):
            font_name = font_name.replace(suffix, "")
        font.name = font_name.strip() or self._default_font

        # Font size
        size = span.size if span.size > 0 else self._default_size
        # Heading overrides
        if style == "heading1":
            size = max(size, 24)
        elif style == "heading2":
            size = max(size, 20)
        elif style == "heading3":
            size = max(size, 16)
        font.size = Pt(size)

        # Bold / italic / underline
        font.bold = span.bold or style in ("heading1", "heading2", "heading3", "heading4")
        font.italic = span.italic
        font.underline = span.underline

        # Color
        if span.color != 0:
            r, g, b = span.color_rgb
            font.color.rgb = RGBColor(r, g, b)

    def _set_alignment(self, pptx_para: Any, alignment: str) -> None:
        """Set paragraph alignment."""
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        pptx_para.alignment = alignment_map.get(alignment, PP_ALIGN.LEFT)

    # ------------------------------------------------------------------
    # Internal — tables
    # ------------------------------------------------------------------

    def _add_table_shape(
        self,
        slide: Any,
        table: FormattedTable,
        left: Any,
        top: Any,
        width: Any,
    ) -> Any:
        """Add a native table shape to the slide. Returns updated y position."""
        if table.is_empty or not table.rows:
            return top

        num_rows = table.num_rows
        num_cols = table.num_cols

        if num_rows == 0 or num_cols == 0:
            return top

        row_height = Pt(20)
        table_height = row_height * num_rows
        col_width = int(width / num_cols) if num_cols > 0 else width

        shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
        pptx_table = shape.table

        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row):
                if c_idx >= num_cols or r_idx >= num_rows:
                    continue

                pptx_cell = pptx_table.cell(r_idx, c_idx)
                pptx_cell.text = cell.text or ""

                # Apply basic formatting
                try:
                    for para in pptx_cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(_TABLE_FONT_SIZE)
                            run.font.name = self._default_font
                            if cell.bold or cell.is_header:
                                run.font.bold = True
                except Exception:
                    pass

        return top + table_height

    # ------------------------------------------------------------------
    # Internal — helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _save(prs: Any) -> bytes:
        """Save a Presentation to bytes."""
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
