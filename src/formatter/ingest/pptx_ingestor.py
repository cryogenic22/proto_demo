"""
PPTX Ingestor — converts a PowerPoint (.pptx) file to FormattedDocument IR.

Uses python-pptx to parse slides, extracting text with formatting metadata
(bold, italic, underline, font, size, color) and tables. Each slide maps
to a FormattedPage; text shapes produce FormattedParagraphs and table shapes
produce FormattedTables.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from src.formatter.extractor import (
    FormattedDocument,
    FormattedLine,
    FormattedPage,
    FormattedParagraph,
    FormattedSpan,
    FormattedTable,
    FormattedTableCell,
)

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

# Default slide dimensions (standard 10" x 7.5" in points)
_DEFAULT_WIDTH = 720.0
_DEFAULT_HEIGHT = 540.0


class PPTXIngestor:
    """Parses a PPTX file into a FormattedDocument IR.

    Usage::

        ingestor = PPTXIngestor()
        doc = ingestor.ingest(pptx_bytes, filename="slides.pptx")
    """

    def ingest(self, content: bytes, filename: str = "") -> FormattedDocument:
        """Parse a PPTX file and return a FormattedDocument.

        Args:
            content: Raw bytes of the .pptx file.
            filename: Original filename for metadata.

        Returns:
            A FormattedDocument populated from the presentation.

        Raises:
            ImportError: If python-pptx is not installed.
        """
        if not _HAS_PPTX:
            raise ImportError(
                "python-pptx is required for PPTX ingestion. "
                "Install it with: pip install python-pptx"
            )

        if not content:
            return FormattedDocument(filename=filename)

        prs = Presentation(io.BytesIO(content))
        pages: list[FormattedPage] = []
        font_counts: dict[str, int] = {}
        color_counts: dict[str, int] = {}
        style_counts: dict[str, int] = {}

        # Slide dimensions
        slide_width = _emu_to_pt(prs.slide_width) if prs.slide_width else _DEFAULT_WIDTH
        slide_height = _emu_to_pt(prs.slide_height) if prs.slide_height else _DEFAULT_HEIGHT

        for slide_idx, slide in enumerate(prs.slides):
            page = FormattedPage(
                page_number=slide_idx,
                width=slide_width,
                height=slide_height,
                margin_left=0,
                margin_right=0,
                margin_top=0,
                margin_bottom=0,
            )

            for shape in slide.shapes:
                if shape.has_table:
                    table = self._extract_table(shape, font_counts, color_counts)
                    if table is not None:
                        page.tables.append(table)
                elif shape.has_text_frame:
                    paragraphs = self._extract_text_frame(
                        shape, font_counts, color_counts, style_counts,
                    )
                    page.paragraphs.extend(paragraphs)

            pages.append(page)

        return FormattedDocument(
            filename=filename,
            pages=pages,
            font_inventory=font_counts,
            color_inventory=color_counts,
            style_inventory=style_counts,
        )

    # ------------------------------------------------------------------
    # Internal — text extraction
    # ------------------------------------------------------------------

    def _extract_text_frame(
        self,
        shape: Any,
        font_counts: dict[str, int],
        color_counts: dict[str, int],
        style_counts: dict[str, int],
    ) -> list[FormattedParagraph]:
        """Extract paragraphs from a shape's text frame."""
        paragraphs: list[FormattedParagraph] = []

        # Shape position for span coordinates
        left = _emu_to_pt(shape.left) if shape.left is not None else 0.0
        top = _emu_to_pt(shape.top) if shape.top is not None else 0.0
        width = _emu_to_pt(shape.width) if shape.width is not None else 0.0

        try:
            text_frame = shape.text_frame
        except Exception:
            return paragraphs

        for pptx_para in text_frame.paragraphs:
            spans: list[FormattedSpan] = []
            x_cursor = left

            for run in pptx_para.runs:
                text = run.text
                if not text:
                    continue

                font = run.font
                font_name = font.name or ""
                font_size = _pt_value(font.size) if font.size is not None else 11.0
                is_bold = bool(font.bold)
                is_italic = bool(font.italic)
                is_underline = bool(font.underline)
                color_int = _rgb_to_int(font.color)

                # Approximate span width
                span_width = len(text) * font_size * 0.5

                span = FormattedSpan(
                    text=text,
                    x0=x_cursor,
                    y0=top,
                    x1=x_cursor + span_width,
                    y1=top + font_size,
                    font=font_name,
                    size=font_size,
                    color=color_int,
                    bold=is_bold,
                    italic=is_italic,
                    underline=is_underline,
                )
                spans.append(span)
                x_cursor += span_width

                # Track inventories
                if font_name:
                    font_counts[font_name] = font_counts.get(font_name, 0) + 1
                hex_color = f"#{color_int:06X}"
                color_counts[hex_color] = color_counts.get(hex_color, 0) + 1

            if not spans:
                continue

            line = FormattedLine(
                spans=spans,
                y_center=top,
                indent=left,
            )

            # Determine style from paragraph level and formatting
            style = self._classify_pptx_paragraph(pptx_para, spans)
            style_counts[style] = style_counts.get(style, 0) + 1

            # Determine alignment
            alignment = _pptx_alignment(pptx_para.alignment)

            para = FormattedParagraph(
                lines=[line],
                style=style,
                indent_level=pptx_para.level if pptx_para.level else 0,
                alignment=alignment,
            )
            paragraphs.append(para)

        return paragraphs

    def _classify_pptx_paragraph(
        self,
        pptx_para: Any,
        spans: list[FormattedSpan],
    ) -> str:
        """Classify a PPTX paragraph into an IR style."""
        if not spans:
            return "body"

        # Check font size of first span
        font_size = spans[0].size
        all_bold = all(s.bold for s in spans if s.text.strip())

        if font_size >= 20 and all_bold:
            return "heading1"
        if font_size >= 16 and all_bold:
            return "heading2"
        if font_size >= 14 and all_bold:
            return "heading3"

        # Bullet detection via level attribute
        level = pptx_para.level if pptx_para.level else 0
        if level > 0:
            return "list_bullet"

        text = "".join(s.text for s in spans).strip()
        if text.startswith(("\u2022", "\u2023", "\u25CF", "\u25CB", "-", "\u2013")):
            return "list_bullet"

        return "body"

    # ------------------------------------------------------------------
    # Internal — table extraction
    # ------------------------------------------------------------------

    def _extract_table(
        self,
        shape: Any,
        font_counts: dict[str, int],
        color_counts: dict[str, int],
    ) -> FormattedTable | None:
        """Extract a FormattedTable from a PPTX table shape."""
        try:
            pptx_table = shape.table
        except Exception:
            return None

        num_rows = len(pptx_table.rows)
        num_cols = len(pptx_table.columns)

        if num_rows == 0 or num_cols == 0:
            return None

        y_position = _emu_to_pt(shape.top) if shape.top is not None else 0.0

        rows: list[list[FormattedTableCell]] = []
        for r_idx, row in enumerate(pptx_table.rows):
            cells: list[FormattedTableCell] = []
            for c_idx, cell in enumerate(row.cells):
                text = cell.text.strip() if cell.text else ""
                is_header = r_idx == 0

                # Check if any run in the cell is bold
                is_bold = is_header
                try:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold:
                                is_bold = True
                            font_name = run.font.name or ""
                            if font_name:
                                font_counts[font_name] = font_counts.get(font_name, 0) + 1
                            color_int = _rgb_to_int(run.font.color)
                            hex_color = f"#{color_int:06X}"
                            color_counts[hex_color] = color_counts.get(hex_color, 0) + 1
                except Exception:
                    pass

                # Detect merged cells via span_height / span_width (if available)
                rowspan = 1
                colspan = 1
                try:
                    if hasattr(cell, "span_height"):
                        rowspan = cell.span_height
                    if hasattr(cell, "span_width"):
                        colspan = cell.span_width
                except Exception:
                    pass

                cells.append(FormattedTableCell(
                    text=text,
                    row=r_idx,
                    col=c_idx,
                    rowspan=rowspan,
                    colspan=colspan,
                    bold=is_bold,
                    is_header=is_header,
                ))
            rows.append(cells)

        return FormattedTable(
            rows=rows,
            num_rows=num_rows,
            num_cols=num_cols,
            y_position=y_position,
        )


# ---------------------------------------------------------------------------
# Module-level helpers
# ---------------------------------------------------------------------------

def _emu_to_pt(emu: int | None) -> float:
    """Convert EMU (English Metric Units) to points."""
    if emu is None:
        return 0.0
    return emu / 12700.0


def _pt_value(emu_size: Any) -> float:
    """Extract a point value from a pptx Pt-like size object."""
    if emu_size is None:
        return 11.0
    try:
        # pptx stores sizes in EMU; Pt objects have a .pt attribute
        if hasattr(emu_size, "pt"):
            return float(emu_size.pt)
        # Raw EMU integer
        return int(emu_size) / 12700.0
    except (TypeError, ValueError):
        return 11.0


def _rgb_to_int(color_obj: Any) -> int:
    """Convert a pptx font color to a packed RGB integer."""
    try:
        if color_obj is None:
            return 0
        rgb = color_obj.rgb
        if rgb is None:
            return 0
        # RGBColor is a tuple-like (r, g, b) with 0-255 values
        return (int(rgb[0]) << 16) | (int(rgb[1]) << 8) | int(rgb[2])
    except Exception:
        return 0


def _pptx_alignment(alignment: Any) -> str:
    """Map pptx alignment enum to IR alignment string."""
    if alignment is None:
        return "left"
    try:
        if _HAS_PPTX:
            mapping = {
                PP_ALIGN.LEFT: "left",
                PP_ALIGN.CENTER: "center",
                PP_ALIGN.RIGHT: "right",
                PP_ALIGN.JUSTIFY: "justify",
            }
            return mapping.get(alignment, "left")
    except Exception:
        pass
    return "left"
